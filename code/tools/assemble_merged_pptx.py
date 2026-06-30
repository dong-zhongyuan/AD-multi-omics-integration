#!/usr/bin/env python3
"""Assemble 3 MERGED main figures from the standardized R panels.

Merge scheme (7 -> 3 figures), each panel keeps its fixed ratio (1:1/2:1/3:1)
and every row tiles to a common width so the grid is tight:

  Fig1 = Neural ODE framework + sensitivity networks (10 panels)
           old fig1(a-f) + fig2(a-d) relabeled a-j, stacked (all rows width 4)
  Fig2 = Hub identification + dual-method KO validation (10 panels)
           old fig3(a-f) + fig4(a-d), relabeled a-j
  Fig3 = Diagnostic + clinical + external validation (9 panels, 3x3 grid)
           old fig5(a-d) + fig6(a-d) + fig7a(AUC heatmap), relabeled a-i

Rows are chosen so each row's panel widths sum to a common total.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor

SRC = os.path.join(os.path.dirname(__file__), '..', 'figures', 'r_figures')
OUT = os.path.join(os.path.dirname(__file__), '..', 'figures',
                   'main_figures_assembly')

SLIDE_W, SLIDE_H = 13.333, 7.5
GAP = 0.10
LABEL_COLOR = RGBColor(0x1A, 0x1A, 0x1A)

# aspect ratio (width : height in height-units) per panel
RATIO = {f'fig{i}{c}': r for i, cs in [
    (1, {'a': 1, 'b': 3, 'c': 1, 'd': 3, 'e': 2, 'f': 2}),
    (2, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (3, {'a': 1, 'b': 2, 'c': 1, 'd': 1, 'e': 1, 'f': 2}),
    (4, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (5, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (6, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (7, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
] for c, r in cs.items()}

# each figure = list of rows; each row = list of (panel_id, label_letter).
# row widths must be equal within a figure.
FIGURES = {
    # Fig1 = Fig1(a-f) + Fig2(a-d)->g-j, stacked; every row width 4
    1: [[('fig1a','a'),('fig1b','b')],                # 1+3=4
        [('fig1c','c'),('fig1d','d')],                # 1+3=4
        [('fig1e','e'),('fig1f','f')],                # 2+2=4
        [('fig2a','g'),('fig2b','h')],                # 2+2=4
        [('fig2c','i'),('fig2d','j')]],               # 2+2=4
    # Fig2 merged: old fig3(a-f) a-f + fig4(a-d) g-j; rows all width 4
    2: [[('fig3a','a'),('fig3b','b'),('fig3c','c')],   # 1+2+1=4
        [('fig3d','d'),('fig3e','e'),('fig3f','f')],   # 1+1+2=4
        [('fig4a','g'),('fig4b','h')],                 # 2+2=4
        [('fig4c','i'),('fig4d','j')]],                # 2+2=4
    # Fig3 = old fig5(a-d) + fig6(a-d) + fig7a(AUC heatmap), 3x3 grid, all 2:1
    3: [[('fig5a','a'),('fig5b','b'),('fig5c','c')],   # 2+2+2=6
        [('fig5d','d'),('fig6a','e'),('fig6b','f')],   # 2+2+2=6
        [('fig6c','g'),('fig6d','h'),('fig7a','i')]],  # 2+2+2=6
}


def panel_path(panel):
    return os.path.join(SRC, 'fig' + panel[3], panel + '.png')


def place(slide, panel, left, top, w, h, label, U):
    slide.shapes.add_picture(panel_path(panel), Inches(left), Inches(top),
                             Inches(w), Inches(h))
    # Label font + box scale with U (the per-figure unit height) so the label
    # is a CONSTANT fraction of each panel regardless of how much the figure
    # was scaled to fit the slide. Without this, figures with more rows (smaller
    # panels, e.g. Fig1) get disproportionately large labels.
    label_in = 0.17 * U
    box_w = 2.0 * label_in
    box_h = 1.6 * label_in
    tb = slide.shapes.add_textbox(Inches(left + 0.06), Inches(top + 0.02),
                                  Inches(box_w), Inches(box_h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Emu(Inches(label_in))
    run.font.bold = True
    run.font.color.rgb = LABEL_COLOR


def build(slide, rows):
    # U scaled to the WIDEST row; each row centered independently on the slide
    row_widths = [sum(RATIO[p] for p, _ in row) for row in rows]
    max_w = max(row_widths)
    n_rows = len(rows)
    U = min((SLIDE_W - 0.6) / (max_w + 0.3), (SLIDE_H - 0.6) / (n_rows + 0.4))
    y0 = (SLIDE_H - n_rows * U - (n_rows - 1) * GAP) / 2 + 0.1
    for ri, row in enumerate(rows):
        rw = row_widths[ri]
        top = y0 + ri * (U + GAP)
        x = (SLIDE_W - rw * U - GAP * (len(row) - 1)) / 2
        for pi, (panel, label) in enumerate(row):
            if pi > 0:
                x += GAP
            cw = RATIO[panel] * U - GAP
            ch = U - GAP if ri < n_rows - 1 else U
            place(slide, panel, x, top, cw, ch, label, U)
            x += RATIO[panel] * U


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    for fig_num in [1, 2, 3]:
        slide = prs.slides.add_slide(blank)
        build(slide, FIGURES[fig_num])
        print(f'Figure {fig_num}: {sum(len(r) for r in FIGURES[fig_num])} panels')
    out_path = os.path.join(OUT, 'figures_merged_5.pptx')
    prs.save(out_path)
    print(f'\nSaved: {out_path}')


if __name__ == '__main__':
    main()
