#!/usr/bin/env python3
"""Assemble R-generated panels into a PPTX with the FIXED original grid.
Panels now have standardized ratios (1:1 / 2:1 / 3:1), so the grid tiles
perfectly: every row is 4 width-units wide (2:1=2, 1:1=1, 3:1=3); fig6 is
5 wide x 2 tall. Letter labels sit inside each panel's top-left corner.
fig6: left col rows [a,b][c], right col d (full height) — user keeps manual.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from PIL import Image

SRC = os.path.join(os.path.dirname(__file__), '..', 'figures', 'r_figures')
OUT = os.path.join(os.path.dirname(__file__), '..', 'figures',
                   'main_figures_assembly')

SLIDE_W, SLIDE_H = 13.333, 7.5
GAP = 0.10
LABEL_COLOR = RGBColor(0x1A, 0x1A, 0x1A)

# fixed aspect ratio (width : height, in height-units) per panel
RATIO = {f'fig{i}{c}': r for i, cs in [
    (1, {'a': 1, 'b': 3, 'c': 1, 'd': 3, 'e': 2, 'f': 2}),
    (2, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (3, {'a': 1, 'b': 2, 'c': 1, 'd': 1, 'e': 1, 'f': 2}),
    (4, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (5, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (6, {'a': 1, 'b': 2, 'c': 3, 'd': 1}),
    (7, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
] for c, r in cs.items()}

ROWS = {
    1: [['fig1a', 'fig1b'], ['fig1c', 'fig1d'], ['fig1e', 'fig1f']],
    2: [['fig2a', 'fig2b'], ['fig2c', 'fig2d']],
    3: [['fig3a', 'fig3b', 'fig3c'], ['fig3d', 'fig3e', 'fig3f']],
    4: [['fig4a', 'fig4b'], ['fig4c', 'fig4d']],
    5: [['fig5a', 'fig5b'], ['fig5c', 'fig5d']],
    7: [['fig7a', 'fig7b'], ['fig7c', 'fig7d']],
}


def panel_path(panel):
    return os.path.join(SRC, 'fig' + panel[3], panel + '.png')


def place(slide, panel, left, top, w, h, label):
    slide.shapes.add_picture(panel_path(panel), Inches(left), Inches(top),
                             Inches(w), Inches(h))
    tb = slide.shapes.add_textbox(Inches(left + 0.06), Inches(top + 0.02),
                                  Inches(0.6), Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Emu(Inches(0.32))
    run.font.bold = True
    run.font.color.rgb = LABEL_COLOR


def build_rows(slide, rows):
    """Every row sums to 4 width-units; 1 unit = U inches tall."""
    total_h = len(rows)
    U = min((SLIDE_W - 0.6) / 4.3, (SLIDE_H - 0.6) / (total_h + 0.4))
    gw = 4 * U
    x0 = (SLIDE_W - gw) / 2
    y0 = (SLIDE_H - total_h * U - (total_h - 1) * GAP) / 2 + 0.1
    for ri, row in enumerate(rows):
        top = y0 + ri * (U + GAP)
        x = x0
        for pi, panel in enumerate(row):
            if pi > 0:
                x += GAP
            cw = RATIO[panel] * U - GAP
            ch = U - GAP if ri < len(rows) - 1 else U
            place(slide, panel, x, top, cw, ch, panel[-1])
            x += RATIO[panel] * U


def build_fig6(slide):
    """left col rows [a,b][c] (3 wide x 2 tall) + right col d (2 wide x 2 tall)
    = 5 wide x 2 tall."""
    U = min((SLIDE_W - 0.6) / 5.3, (SLIDE_H - 0.6) / 2.4)
    gw, gh = 5 * U, 2 * U
    x0 = (SLIDE_W - gw) / 2
    y0 = (SLIDE_H - gh) / 2 + 0.1
    place(slide, 'fig6a', x0, y0, 1 * U - GAP, U - GAP, 'a')
    place(slide, 'fig6b', x0 + 1 * U, y0, 2 * U - GAP, U - GAP, 'b')
    place(slide, 'fig6c', x0, y0 + U, 3 * U - GAP, U, 'c')
    place(slide, 'fig6d', x0 + 3 * U, y0, 2 * U - GAP, 2 * U - GAP, 'd')


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    for fig_num in [1, 2, 3, 4, 5, 6, 7]:
        slide = prs.slides.add_slide(blank)
        if fig_num == 6:
            build_fig6(slide)
        else:
            build_rows(slide, ROWS[fig_num])
        print(f'Figure {fig_num}: slide built')
    out_path = os.path.join(OUT, 'figures_r_compiled.pptx')
    prs.save(out_path)
    print(f'\nSaved: {out_path}')


if __name__ == '__main__':
    main()
