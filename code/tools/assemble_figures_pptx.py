#!/usr/bin/env python3
"""Assemble per-figure panels into a PPTX, one slide per figure, with a/b/c...
letter labels. Panels are tiled on a grid where one "height unit" = U inches,
and each panel's cell is sized by its aspect ratio so rows/columns meet
edge-to-edge with no gaps.

Grid layouts (panel -> cell in units of U):
  fig1: rows (a,b)(c,d)(e,f)            -> 4 wide x 3 tall
  fig2: rows (a,b)(c,d)                 -> 4 wide x 2 tall
  fig3: rows (a,b,c)(d,e,f)             -> 4 wide x 2 tall
  fig4: rows (a,b)(c,d)                 -> 4 wide x 2 tall
  fig5: rows (a,b)(c,d)                 -> 4 wide x 2 tall
  fig6: left col rows (a,b)(c) | right col d  -> 5 wide x 2 tall
  fig7: rows (a,b)(c,d)                 -> 4 wide x 2 tall
"""
import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from PIL import Image

PANEL_DIR = os.path.join(os.path.dirname(__file__), '..', 'figures', '_panels')
OUT = os.path.join(os.path.dirname(__file__), '..', 'figures',
                   'main_figures_assembly')

# slide: 13.333 x 7.5 inches (16:9)
SLIDE_W = 13.333
SLIDE_H = 7.5
LABEL_FONT = 28
LABEL_COLOR = RGBColor(0x1A, 0x1A, 0x1A)

# aspect ratio (w:h in height-units) of each panel, in placement order
RATIO = {f'fig{i}{c}': r for i, cs in [
    (1, {'a': 1, 'b': 3, 'c': 1, 'd': 3, 'e': 2, 'f': 2}),
    (2, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (3, {'a': 1, 'b': 2, 'c': 1, 'd': 1, 'e': 1, 'f': 2}),
    (4, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (5, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
    (6, {'a': 1, 'b': 2, 'c': 3, 'd': 1}),
    (7, {'a': 2, 'b': 2, 'c': 2, 'd': 2}),
] for c, r in cs.items()}

# each figure -> list of rows; each row -> list of panel ids placed left->right
# at a shared row height. fig6 is special (see build_fig6).
ROWS = {
    1: [['fig1a', 'fig1b'], ['fig1c', 'fig1d'], ['fig1e', 'fig1f']],
    2: [['fig2a', 'fig2b'], ['fig2c', 'fig2d']],
    3: [['fig3a', 'fig3b', 'fig3c'], ['fig3d', 'fig3e', 'fig3f']],
    4: [['fig4a', 'fig4b'], ['fig4c', 'fig4d']],
    5: [['fig5a', 'fig5b'], ['fig5c', 'fig5d']],
    7: [['fig7a', 'fig7b'], ['fig7c', 'fig7d']],
}


def place_panel(slide, panel, left, top, cell_w, cell_h, label):
    """Place a panel image fit into (cell_w, cell_h) preserving aspect ratio,
    anchored top-left within the cell. The letter label sits just inside the
    panel's own top-left corner (overlapping the panel margin, not its data)."""
    img_path = os.path.join(PANEL_DIR, f'{panel}.png')
    iw, ih = Image.open(img_path).size
    ar = iw / ih
    # fit inside the cell, anchored top-left
    w_by_h = cell_h * ar
    h_by_w = cell_w / ar
    if w_by_h <= cell_w:
        w, h = w_by_h, cell_h
    else:
        w, h = cell_w, h_by_w
    slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                             Inches(w), Inches(h))
    # letter label just inside the panel's top-left corner
    tb = slide.shapes.add_textbox(Inches(left + 0.06), Inches(top + 0.02),
                                  Inches(0.6), Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = label
    p.runs[0].font.size = Emu(Inches(0.34))
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = LABEL_COLOR


# small gap between panels (inches), so adjacent panels don't touch and the
# letter labels land on each panel's own top-left corner
GAP = 0.12


def build_rows_figure(slide, rows, fig_num):
    """Generic 4-wide grid: every row has total width 4 units, with a small
    gap between adjacent panels (horizontal within a row, vertical between
    rows). The gap is subtracted so the outer grid keeps its size."""
    total_h = len(rows)            # height units
    U = min(SLIDE_W / 4.2, SLIDE_H / (total_h + 0.4))  # inch per unit
    grid_w = 4 * U
    grid_h = total_h * U
    x0 = (SLIDE_W - grid_w) / 2
    y0 = (SLIDE_H - grid_h) / 2 + 0.15
    for ri, row in enumerate(rows):
        row_top = y0 + ri * U + ri * GAP     # shift down by accumulated row gaps
        x = x0
        for pi, panel in enumerate(row):
            if pi > 0:
                x += GAP                      # gap between horizontal panels
            cw = RATIO[panel] * U - GAP       # shrink cell to make room for gap
            ch = U - GAP if ri < len(rows) - 1 else U  # leave room for next row gap
            letter = panel[-1]
            place_panel(slide, panel, x, row_top, cw, ch, letter)
            x += RATIO[panel] * U             # advance by full nominal width


def build_fig6(slide):
    """Left half (3 wide x 2 tall): top row [a(1), b(2)], bottom [c(3)].
    Right half: d spanning full height, width 2 -> total 5 wide x 2 tall."""
    U = min(SLIDE_W / 5.2, SLIDE_H / 2.4)
    grid_w = 5 * U
    grid_h = 2 * U
    x0 = (SLIDE_W - grid_w) / 2
    y0 = (SLIDE_H - grid_h) / 2 + 0.15
    half = U  # each of the 2 rows is 1 unit tall
    # left top: a(1U wide) + b(2U wide)
    place_panel(slide, 'fig6a', x0, y0, 1 * U, half, 'a')
    place_panel(slide, 'fig6b', x0 + 1 * U, y0, 2 * U, half, 'b')
    # left bottom: c (3U wide)
    place_panel(slide, 'fig6c', x0, y0 + half, 3 * U, half, 'c')
    # right: d (2U wide, full height 2U)
    place_panel(slide, 'fig6d', x0 + 3 * U, y0, 2 * U, 2 * U, 'd')


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    for fig_num in [1, 2, 3, 4, 5, 6, 7]:
        slide = prs.slides.add_slide(blank)
        if fig_num == 6:
            build_fig6(slide)
        else:
            build_rows_figure(slide, ROWS[fig_num], fig_num)
        print(f'Figure {fig_num}: slide built')

    out_path = os.path.join(OUT, 'figures_compiled.pptx')
    prs.save(out_path)
    print(f'\nSaved: {out_path}')


if __name__ == '__main__':
    main()
