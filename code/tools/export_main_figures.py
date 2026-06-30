#!/usr/bin/env python3
"""Export each PPTX slide as a tight high-resolution main figure.

Reads the user-adjusted PPTX, and for each slide:
  1. extracts the embedded picture blobs (their native pixel images),
  2. uses the slide's real shape coordinates (inches) to lay them out,
  3. draws the a/b/c... letter labels at their real text-box positions,
  4. crops to the bounding box of all shapes (no slide-margin whitespace),
  5. renders at high DPI and saves PNG + TIFF.

Output is the figure content only, sized to its true proportions.
"""
import os
import io
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

PPTX = os.path.join(os.path.dirname(__file__), '..', 'figures',
                    'main_figures_assembly', 'figures_compiled.pptx')
OUT = os.path.join(os.path.dirname(__file__), '..', 'figures',
                   'main_figures_assembly', 'exported')
DPI = 300          # px per inch in the output
LABEL_PT = 26      # label font size (pt) — matches the pptx 0.34in ~ 24pt scale


def inch(v):
    return v / 914400.0 if v is not None else 0.0


def slide_shapes(slide):
    """Return list of dicts: pictures and textboxes with inch coords."""
    pics, texts = [], []
    for sh in slide.shapes:
        if sh.shape_type == 13:  # PICTURE
            blob = sh.image.blob
            pics.append({
                'img': Image.open(io.BytesIO(blob)).convert('RGB'),
                'L': inch(sh.left), 'T': inch(sh.top),
                'W': inch(sh.width), 'H': inch(sh.height),
            })
        elif sh.has_text_frame and sh.text_frame.text.strip():
            texts.append({
                'text': sh.text_frame.text.strip(),
                'L': inch(sh.left), 'T': inch(sh.top),
            })
    return pics, texts


def render(pics, texts):
    # bounding box of all shapes (pictures + text labels)
    min_x = min([p['L'] for p in pics] + [t['L'] for t in texts])
    min_y = min([p['T'] for p in pics] + [t['T'] for t in texts])
    max_x = max([p['L'] + p['W'] for p in pics] +
                [t['L'] + 0.4 for t in texts])
    max_y = max([p['T'] + p['H'] for p in pics] +
                [t['T'] + 0.35 for t in texts])
    W = max_x - min_x
    H = max_y - min_y

    canvas = Image.new('RGB', (int(W * DPI), int(H * DPI)), 'white')
    # place each picture fit into its target box (stretch to box = how pptx shows it)
    for p in pics:
        iw, ih = p['img'].size
        tw, th = int(p['W'] * DPI), int(p['H'] * DPI)
        resized = p['img'].resize((tw, th), Image.LANCZOS)
        canvas.paste(resized, (int((p['L'] - min_x) * DPI),
                               int((p['T'] - min_y) * DPI)))
    # labels
    try:
        font = ImageFont.truetype("arialbd.ttf", int(LABEL_PT * DPI / 72))
    except Exception:
        font = ImageFont.load_default()
    d = ImageDraw.Draw(canvas)
    for t in texts:
        x = int((t['L'] - min_x) * DPI)
        y = int((t['T'] - min_y) * DPI)
        d.text((x, y), t['text'], fill='black', font=font)
    return canvas


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation(PPTX)
    for i, slide in enumerate(prs.slides, 1):
        pics, texts = slide_shapes(slide)
        if not pics:
            continue
        canvas = render(pics, texts)
        name = f'figure{i}'
        png = os.path.join(OUT, f'{name}.png')
        canvas.save(png, dpi=(DPI, DPI))
        canvas.save(os.path.join(OUT, f'{name}.tiff'), dpi=(DPI, DPI))
        print(f'{name}: {canvas.size[0]}x{canvas.size[1]} px  '
              f'(~{canvas.size[0]/DPI:.1f}x{canvas.size[1]/DPI:.1f} in)')


if __name__ == '__main__':
    main()
